import json
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pdfplumber  # 记得添加这个导入

def read_pdf(pdf_path):
    """读取PDF文件并返回文本内容"""
    full_text = ""
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                full_text += text + "\n"
    return full_text

def chunk_text(text, max_chars=3000):
    """将长文本按最大字符数分块"""
    chunks = []
    for i in range(0, len(text), max_chars):
        chunks.append(text[i:i + max_chars])
    return chunks

# ==================== 1. API 配置 ====================
from dotenv import load_dotenv
import os
from openai import OpenAI   # 确保这行在文件顶部已经有了

load_dotenv()
api_key = os.getenv("DEEPSEEK_API_KEY")
if not api_key:
    raise ValueError("请在 .env 文件中设置 DEEPSEEK_API_KEY")

client = OpenAI(
    api_key=api_key,
    base_url="https://api.deepseek.com"
)


# ==================== 2. AI 调用 ====================
def call_ai(user_content, system_content=None):
    messages = []
    if system_content:
        messages.append({"role": "system", "content": system_content})
    messages.append({"role": "user", "content": user_content})

    response = client.chat.completions.create(
        model="deepseek-chat",
        messages=messages,
        response_format={'type': 'json_object'},
        stream=False
    )
    return response.choices[0].message.content


# ==================== 3. PPT 生成 ====================
def create_ppt(data):

    # 兜底检查（保护整个函数）
    if "three_levels" not in data or len(data["three_levels"]) != 3:
        print("⚠️ 三大层级数据异常，使用默认值")
        data["three_levels"] = [
            {"level": "Hardware", "features": "硬件技术层（默认值）"},
            {"level": "Algorithm", "features": "算法与模型层（默认值）"},
            {"level": "Application", "features": "应用与社会层（默认值）"}
        ]
    if "timeline" not in data:
        data["timeline"] = {
            "near_term_2_5_years": ["默认技术1", "默认技术2"],
            "far_term_6_10_years": ["默认技术3", "默认技术4"]
        }

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # =============================================
    # 第一页：标题 + 核心愿景（完全手动控制）
    # =============================================
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    # 主标题（居中）
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11.333), Inches(1.2))
    tf = title_box.text_frame
    tf.text = data.get("title", "AI+硬件协同设计愿景简报")
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 副标题（居中）
    sub_box = slide1.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.333), Inches(0.8))
    tf = sub_box.text_frame
    tf.text = "基于论文：AI+HV2035: Shaping the Next Decade (arXiv:2603.05225)"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 核心愿景（居中）
    core_text = data.get("core_goal", "未提取到核心目标")
    core_box = slide1.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(10.333), Inches(2.5))
    tf = core_box.text_frame
    tf.word_wrap = True
    tf.text = f"🎯 核心目标：\n{core_text}"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    if len(tf.paragraphs) > 1:
        tf.paragraphs[1].font.size = Pt(18)
        tf.paragraphs[1].font.bold = False
        tf.paragraphs[1].alignment = PP_ALIGN.CENTER

    # 底部脚注
    footer_box = slide1.shapes.add_textbox(Inches(1), Inches(6.8), Inches(11.333), Inches(0.4))
    tf = footer_box.text_frame
    tf.text = "AI+硬件协同设计愿景 | 自动化简报生成器"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # =============================================
    # 第二页：三大层级（表格版）
    # =============================================
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title2.text_frame
    tf.text = "三大协同设计层级"
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    table = slide2.shapes.add_table(4, 2, Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5)).table
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(6.833)

    # 表头（蓝底白字）
    for col_idx in range(2):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(18)
    table.cell(0, 0).text = "抽象层级"
    table.cell(0, 1).text = "关键特征"

    if "three_levels" in data and len(data["three_levels"]) >= 3:
        for i, level_info in enumerate(data["three_levels"], start=1):
            table.cell(i, 0).text = level_info.get("level", f"层级{i}")
            table.cell(i, 1).text = level_info.get("features", "")
            for col_idx in range(2):
                cell = table.cell(i, col_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
                cell.text_frame.paragraphs[0].font.size = Pt(16)
    else:
        table.cell(1, 0).text = "未提取到数据"
        table.cell(1, 1).text = "请检查论文内容或提示词"

    # =============================================
    # 第三页：时间轴（干净版本，无多余元素）
    # =============================================
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title3.text_frame
    tf.text = "技术发展时间轴"
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 近期（左）
    near_box = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4.5), Inches(4.0))
    tf = near_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "📌 近期（2-5年）"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    for item in data.get("timeline", {}).get("near_term_2_5_years", []):
        p = tf.add_paragraph()
        p.text = f"  • {item}"
        p.font.size = Pt(16)
        p.space_after = Pt(8)

    # 中间分隔线（代替圆柱）
    line = slide3.shapes.add_shape(1, Inches(5.8), Inches(2.0), Inches(0.04), Inches(3.5))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(200, 200, 200)
    line.line.color.rgb = RGBColor(200, 200, 200)

    # 远期（右）
    far_box = slide3.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(4.5), Inches(4.0))
    tf = far_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "🚀 远期（6-10年）"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(237, 125, 49)
    for item in data.get("timeline", {}).get("far_term_6_10_years", []):
        p = tf.add_paragraph()
        p.text = f"  • {item}"
        p.font.size = Pt(16)
        p.space_after = Pt(8)

    # 底部时间轴横线（视觉美化）
    timeline_line = slide3.shapes.add_shape(1, Inches(0.8), Inches(6.2), Inches(11.733), Inches(0.04))
    timeline_line.fill.solid()
    timeline_line.fill.fore_color.rgb = RGBColor(200, 200, 200)
    timeline_line.line.color.rgb = RGBColor(200, 200, 200)

    # 添加两个时间节点标记
    left_dot = slide3.shapes.add_shape(1, Inches(0.8), Inches(6.0), Inches(0.3), Inches(0.3))
    left_dot.fill.solid()
    left_dot.fill.fore_color.rgb = RGBColor(0, 112, 192)
    left_dot.line.color.rgb = RGBColor(0, 112, 192)

    right_dot = slide3.shapes.add_shape(1, Inches(12.2), Inches(6.0), Inches(0.3), Inches(0.3))
    right_dot.fill.solid()
    right_dot.fill.fore_color.rgb = RGBColor(237, 125, 49)
    right_dot.line.color.rgb = RGBColor(237, 125, 49)

    prs.save("briefing.pptx")
    print("✅ PPT 已生成：briefing.pptx")

# ==================== 4. 主程序 ====================
if __name__ == "__main__":
    print("=== Day 2: 提取完整结构化信息 ===")

    # 1. 读PDF
    print("正在读论文...")
    full_text = read_pdf("paper.pdf")
    print(f"读到了 {len(full_text)} 字符")

    # 2. 智能定位论文核心章节（精度优化）
    keyword = "2.2"  # 论文第 2.2 节（三大层级的定义处）
    index = full_text.find(keyword)
    if index != -1:
        # 找到了！只取关键词前后共 3000 字符，精准且省 Token
        start = max(0, index - 500)
        end = min(len(full_text), index + 4000)
        first_part = full_text[start:end]
        print(f"✅ 已精准定位到 2.2 章节，截取长度：{len(first_part)} 字符")
    else:
        # 万一没找到（兜底方案），退回原来的分块策略
        print("⚠️ 未找到 2.2 关键词，使用兜底分块方案")
        chunks = chunk_text(full_text, max_chars=3000)
        first_part = chunks[0] + "\n" + (chunks[1] if len(chunks) > 1 else "")
        print(f"喂给AI的文本长度：{len(first_part)} 字符")

    # 3. 构造prompt
    prompt = f"""
    你是一个资深技术分析师。请从以下论文内容中提取结构化信息，并以JSON格式返回。

    论文内容：
    {first_part}

    请提取以下内容：
    1. "title": 论文标题（英文原样）
    2. "core_goal": 用一句中文描述"1000倍效率提升"的定义
    3. "three_levels": 论文2.2节的三个抽象层级（Hardware / Algorithm / Application）及其特征
    4. "timeline": 近期（2-5年）和远期（6-10年）的技术趋势

    请严格按照以下JSON格式返回：
    {{
        "title": "论文标题",
        "core_goal": "1000倍效率提升的定义",
        "three_levels": [
            {{"level": "Hardware", "features": "特征描述"}},
            {{"level": "Algorithm", "features": "特征描述"}},
            {{"level": "Application", "features": "特征描述"}}
        ],
        "timeline": {{
            "near_term_2_5_years": ["技术1", "技术2"],
            "far_term_6_10_years": ["技术1", "技术2"]
        }}
    }}
    """

    # 4. 调用AI
    print("正在呼叫DeepSeek提取完整信息...")
    try:
        reply = call_ai(prompt)
        print("AI原始返回：", reply)

        # 5. 解析JSON
        data = json.loads(reply)

        # 打印漂亮的结果
        print("\n" + "=" * 50)
        print("✅ 提取成功！结果如下：")
        print(json.dumps(data, ensure_ascii=False, indent=2))
        create_ppt(data)
        print("=" * 50)
    except Exception as e:
        print(f"❌ 出错：{e}")
        print("原始返回内容：", reply if 'reply' in locals() else "无返回")